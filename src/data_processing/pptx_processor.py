"""
PPTX processor — extracts text from slides and speaker notes.

Tracks slide_number metadata for citation.
"""

from pathlib import Path
from dataclasses import dataclass, field
from pptx import Presentation

from src.utils.logger import setup_logger

logger = setup_logger(__name__)


@dataclass
class SlideContent:
    """Extracted content from a single PPTX slide."""
    slide_number: int
    title: str = ""
    body_text: str = ""
    speaker_notes: str = ""
    shapes_text: list[str] = field(default_factory=list)
    has_table: bool = False
    table_data: list[list[str]] = field(default_factory=list)


class PPTXProcessor:
    """
    Extracts text content from PowerPoint presentations.
    Captures slide text, titles, speaker notes, and table data.
    """

    def process(self, pptx_path: str, doc_id: str) -> list[SlideContent]:
        """
        Processes a PPTX file into per-slide content objects.

        Args:
            pptx_path: Absolute path to the PPTX file.
            doc_id: Document identifier for metadata.

        Returns:
            List of SlideContent objects, one per slide.

        Raises:
            FileNotFoundError: If pptx_path does not exist.
        """
        path = Path(pptx_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX not found: {pptx_path}")

        logger.info(
            "Processing PPTX",
            extra={"pptx_path": pptx_path, "doc_id": doc_id},
        )

        prs = Presentation(pptx_path)
        slides: list[SlideContent] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            content = SlideContent(slide_number=slide_num)

            # Extract title
            if slide.shapes.title:
                content.title = slide.shapes.title.text.strip()

            # Extract text from all shapes
            body_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            # Skip title (already captured)
                            if shape == slide.shapes.title:
                                continue
                            body_parts.append(text)
                            content.shapes_text.append(text)

                # Extract table data
                if shape.has_table:
                    content.has_table = True
                    table = shape.table
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        content.table_data.append(row_data)

            content.body_text = "\n".join(body_parts)

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_parts = []
                for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        notes_parts.append(text)
                content.speaker_notes = "\n".join(notes_parts)

            slides.append(content)

        logger.info(
            "PPTX processing complete",
            extra={
                "doc_id": doc_id,
                "total_slides": len(slides),
                "slides_with_tables": sum(1 for s in slides if s.has_table),
                "slides_with_notes": sum(1 for s in slides if s.speaker_notes),
            },
        )

        return slides

    def to_text_chunks(self, slides: list[SlideContent]) -> list[dict]:
        """
        Converts SlideContent objects to text chunks for chunking pipeline.

        Args:
            slides: List of SlideContent from process().

        Returns:
            List of dicts with text, slide_number, section_heading, content_type.
        """
        chunks = []
        for slide in slides:
            # Main slide content
            parts = []
            if slide.title:
                parts.append(f"# {slide.title}")
            if slide.body_text:
                parts.append(slide.body_text)

            if parts:
                chunks.append({
                    "text": "\n\n".join(parts),
                    "slide_number": slide.slide_number,
                    "section_heading": slide.title,
                    "content_type": "slide",
                })

            # Speaker notes as separate chunk
            if slide.speaker_notes:
                chunks.append({
                    "text": f"[Speaker Notes - Slide {slide.slide_number}]\n{slide.speaker_notes}",
                    "slide_number": slide.slide_number,
                    "section_heading": f"{slide.title} (Notes)",
                    "content_type": "speaker_notes",
                })

            # Table data as separate chunk
            if slide.table_data:
                table_lines = []
                for row in slide.table_data:
                    table_lines.append(" | ".join(row))
                chunks.append({
                    "text": "\n".join(table_lines),
                    "slide_number": slide.slide_number,
                    "section_heading": f"{slide.title} (Table)",
                    "content_type": "table",
                    "is_table": True,
                })

        return chunks
